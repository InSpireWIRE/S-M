import PyPDF2
import pdfplumber
from pptx import Presentation
from docx import Document
import json
from datetime import datetime
from deck_chunker import UniversalDocumentaryChunker
from advanced_chunker import AdvancedDocumentaryChunker
import uuid
import asyncio
import requests
from url_scraper import AdvancedURLScraper
from typing import Dict

class DeckUploadHandler:
    def __init__(self, supabase_client):
        self.supabase = supabase_client
    
    def extract_pdf_text(self, file_path):
        """Extract text from PDF using pdfplumber (better for tables)"""
        text = ""
        try:
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n\n"
        except:
            # Fallback to PyPDF2
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    text += page.extract_text() + "\n\n"
        return text
    
    def extract_pptx_text(self, file_path):
        """Extract text from PowerPoint"""
        prs = Presentation(file_path)
        text = ""
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
            text += "\n---SLIDE BREAK---\n\n"
        
        return text
    
    def extract_docx_text(self, file_path):
        """Extract text from Word document"""
        doc = Document(file_path)
        text = ""
        
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        
        return text
    
    def process_upload(self, file_path, filename, user_id):
        """Process uploaded file and store in database"""
        
        # Determine file type and extract text
        if filename.lower().endswith('.pdf'):
            raw_text = self.extract_pdf_text(file_path)
            file_type = 'pdf'
        elif filename.lower().endswith(('.ppt', '.pptx')):
            raw_text = self.extract_pptx_text(file_path)
            file_type = 'powerpoint'
        elif filename.lower().endswith('.docx'):
            raw_text = self.extract_docx_text(file_path)
            file_type = 'word'
        elif filename.lower().endswith('.doc'):
            # For testing - treat .doc as plain text
            with open(file_path, 'r') as f:
                raw_text = f.read()
            file_type = 'word'    
        else:
            raise ValueError(f"Unsupported file type: {filename}")
        
        # Store in database
        deck_data = {
            'user_id': 'c50f98ec-1234-5678-9abc-def012345678',
            'deck_name': filename,
            'deck_type': file_type,
            'original_filename': filename,
            'content_extracted': {
                'raw_text': raw_text,
                'upload_time': datetime.now().isoformat(),
                'file_size': len(raw_text),
                'word_count': len(raw_text.split())
            }
        }
        
        result = self.supabase.table('uploaded_decks').insert(deck_data).execute()
        
        # NOW DO THE CHUNKING - INSIDE THE METHOD!
        if result.data:
            deck_id = result.data[0]['id']
            
            # Process deck into intelligent chunks
            chunker = AdvancedDocumentaryChunker(self.supabase)
            chunk_result = chunker.process_deck_advanced(deck_id, raw_text, file_type)
            
            # Update deck with chunking analysis
            update_data = {
                'status': 'ready',
                'content_extracted': {
                    **deck_data['content_extracted'],  # Keep existing content
                    'genre': chunk_result.get('synthesis', {}).get('narrative_flow', 'unknown'),
                    'structure_type': chunk_result.get('synthesis', {}).get('structure_type', 'unknown'),
                    'total_chunks': len(chunk_result.get('chunks', [])),
                    'elements_found': chunk_result.get('synthesis', {}).get('key_entities', {}),
                    'missing_elements': chunk_result.get('synthesis', {}).get('central_themes', [])
                }
            }
            
            self.supabase.table('uploaded_decks').update(update_data).eq('id', deck_id).execute()
            
            return result.data[0]

    def process_url(self, url: str, user_id: str) -> Dict:
        """Process a deck from URL with multi-page support"""
        
        deck_id = str(uuid.uuid4())
        
        try:
            # Use the advanced scraper
            scraper = AdvancedURLScraper()
            
            # Run async scraping
            if 'indd.adobe.com' in url:
                raw_text = asyncio.run(scraper.scrape_adobe_indesign(url))
            else:
                raw_text = asyncio.run(scraper.scrape_generic_presentation(url))
            
            if not raw_text or len(raw_text) < 50:
                return {'error': 'Could not extract sufficient content from URL'}
            
            # Store initial deck record
            deck_data = {
                'id': deck_id,
                'user_id': user_id,
                'deck_name': f"url_{deck_id[:8]}",
                'deck_type': 'word',  # Use 'word' as a valid type for URLs
                'original_filename': f"url_{deck_id[:8]}.txt",
                'content_extracted': {
                    'raw_text': raw_text,
                    'upload_time': datetime.now().isoformat(),
                    'file_size': len(raw_text),
                    'word_count': len(raw_text.split()),  # <-- ADD COMMA HERE
                    'source_url': url  # Track the original URL here
                }
            }
            
            self.supabase.table('uploaded_decks').insert(deck_data).execute()
            
            # Process with advanced chunker
            chunker = AdvancedDocumentaryChunker(self.supabase)
            chunk_result = chunker.process_deck_advanced(deck_id, raw_text, 'url')
            
            # Update deck with analysis
            update_data = {
                'content_extracted': {
                    **deck_data['content_extracted'],
                    'genre': chunk_result.get('synthesis', {}).get('narrative_flow', 'unknown'),
                    'structure_type': chunk_result.get('synthesis', {}).get('structure_type', 'unknown'),
                    'total_chunks': len(chunk_result.get('chunks', [])),
                    'elements_found': chunk_result.get('synthesis', {}).get('key_entities', {}),
                    'missing_elements': chunk_result.get('synthesis', {}).get('central_themes', [])
                }
            }
            
            self.supabase.table('uploaded_decks').update(update_data).eq('id', deck_id).execute()
            
            return {
                'deck_id': deck_id,
                'message': 'URL processed successfully',
                'word_count': len(raw_text.split()),
                'source': url
            }
            
        except Exception as e:
            return {'error': str(e)}

        return None